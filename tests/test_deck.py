"""The slide deck and its exported PDF.

These exist because the PDF is a *derived* artifact that nothing regenerates automatically.
Rebuilding the deck updates the .pptx; the .pdf only changes if someone remembers to export
it again. That silently drifted to a 29-page PDF beside a 30-slide deck, and since the PDF
is what actually gets emailed or printed, the stale one is the one people read.
"""
from __future__ import annotations

import os

import pytest

DECK = os.path.join("Assignment 1", "slides", "Components_1_3_Update.pptx")
PDF = os.path.join("Assignment 1", "slides", "Components_1_3_Update.pdf")


@pytest.fixture(scope="module")
def deck(repo):
    pptx = pytest.importorskip("pptx")
    return pptx.Presentation(os.path.join(repo, DECK))


def test_pdf_export_matches_the_deck(repo, deck):
    """One page per slide. A mismatch means the PDF was exported from an older deck."""
    fitz = pytest.importorskip("pymupdf")
    pdf_path = os.path.join(repo, PDF)
    assert os.path.exists(pdf_path), "the exported PDF is missing"
    pages = fitz.open(pdf_path).page_count
    assert pages == len(deck.slides), (
        f"PDF has {pages} pages but the deck has {len(deck.slides)} slides - "
        "re-export the PDF from the current .pptx")


def test_every_slide_has_a_title(deck):
    """A slide with no text is almost always a build error rather than a design choice."""
    empty = [i for i, s in enumerate(deck.slides, 1)
             if not any(sh.has_text_frame and sh.text_frame.text.strip() for sh in sh_list(s))]
    assert not empty, f"slides with no text at all: {empty}"


def sh_list(slide):
    return list(slide.shapes)


def test_no_slide_overflows_the_footer(deck):
    """The footer sits at y = 7.0 in. Content running under it is unreadable when projected.

    Caught a real overlap once: a takeaway bar grew past the footer after its text was
    lengthened, and nothing about the deck build complained.
    """
    from pptx.util import Emu
    FOOTER_Y = 7.0
    offenders = []
    for i, slide in enumerate(deck.slides, 1):
        for sh in slide.shapes:
            top = Emu(sh.top).inches if sh.top is not None else 0.0
            height = Emu(sh.height).inches if sh.height is not None else 0.0
            if top < FOOTER_Y - 0.02 and top + height > FOOTER_Y + 0.02:
                offenders.append((i, round(top + height, 2)))
    assert not offenders, f"content crossing the footer line at y=7.0: {offenders}"


def test_deck_and_companion_docs_cover_the_same_slides(repo, deck):
    """The speaking notes live outside the repo, so this only runs where they exist.

    Renumbering the deck without renumbering the companion docs has happened repeatedly:
    the docs stay internally consistent while pointing at the wrong slides.
    """
    import io
    import re
    onedrive = os.path.join(
        os.path.expanduser("~"), "OneDrive", "Documents", "Research Internship",
        "Quantum Research Internship", "Assignment 1", "slides")
    if not os.path.isdir(onedrive):
        pytest.skip("companion docs not present on this machine")
    n = len(deck.slides)
    for fn in ("Meeting_Script.md", "Slide_by_Slide.md"):
        path = os.path.join(onedrive, fn)
        if not os.path.exists(path):
            continue
        text = io.open(path, encoding="utf-8").read()
        covered = set()
        for m in re.finditer(r"^## Slides? ([0-9–]+) —", text, re.M):
            rng = m.group(1)
            if "–" in rng:
                a, b = rng.split("–")
                covered |= set(range(int(a), int(b) + 1))
            else:
                covered.add(int(rng))
        missing = sorted(set(range(1, n + 1)) - covered)
        assert not missing, f"{fn} has no block for slides {missing}"
